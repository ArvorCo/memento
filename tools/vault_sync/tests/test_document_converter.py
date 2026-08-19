from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from tools.vault_sync.document_converter import _clean_pdf_pages, convert_document


class DocumentConverterTests(unittest.TestCase):
    def test_csv_conversion_preserves_rows_as_markdown_table(self) -> None:
        with tempfile.TemporaryDirectory() as tempdir:
            path = Path(tempdir) / "people.csv"
            path.write_text('name,role\nAda,"engineer|researcher"\n', encoding="utf-8")

            result = convert_document(path)

            self.assertEqual(result.converter, "python-csv")
            self.assertIn("| name | role |", result.markdown)
            self.assertIn(r"engineer\|researcher", result.markdown)

    def test_pdf_cleanup_removes_repeated_edges_and_dehyphenates_words(self) -> None:
        pages = [
            "Quarterly Report\ninter-\nnational results\n1",
            "Quarterly Report\nsecond page\n2",
            "Quarterly Report\nthird page\n3",
        ]

        cleaned = _clean_pdf_pages(pages)

        self.assertNotIn("Quarterly Report", "\n".join(cleaned))
        self.assertIn("international results", cleaned[0])

    def test_docx_conversion_uses_native_python_reader(self) -> None:
        from docx import Document

        with tempfile.TemporaryDirectory() as tempdir:
            path = Path(tempdir) / "plan.docx"
            document = Document()
            document.add_heading("Windows plan", level=1)
            document.add_paragraph("Named pipes keep the runtime local.")
            document.save(path)

            result = convert_document(path)

            self.assertEqual(result.converter, "python-docx")
            self.assertIn("Windows plan", result.markdown)
            self.assertIn("Named pipes", result.markdown)

    def test_pptx_conversion_preserves_slide_text(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tempdir:
            path = Path(tempdir) / "briefing.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Memento"
            slide.placeholders[1].text = "Local Windows memory"
            presentation.save(path)

            result = convert_document(path)

            self.assertEqual(result.converter, "python-pptx")
            self.assertIn("Local Windows memory", result.markdown)

    def test_xlsx_conversion_preserves_sheets_and_cells(self) -> None:
        from openpyxl import Workbook

        with tempfile.TemporaryDirectory() as tempdir:
            path = Path(tempdir) / "facts.xlsx"
            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "Decisions"
            worksheet.append(["topic", "decision"])
            worksheet.append(["transport", "named pipe"])
            workbook.save(path)

            result = convert_document(path)

            self.assertEqual(result.converter, "openpyxl")
            self.assertIn("## Sheet: Decisions", result.markdown)
            self.assertIn("| transport | named pipe |", result.markdown)
