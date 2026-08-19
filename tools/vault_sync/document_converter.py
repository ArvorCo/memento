from __future__ import annotations

import csv
import html
import importlib.util
import io
import json
import re
import shutil
import subprocess
import tempfile
from collections import Counter
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path

from tools.vault_sync.markdown import normalize_markdown

PANDOC_EXTENSIONS = {
    ".docx",
    ".epub",
    ".html",
    ".htm",
    ".ipynb",
    ".odt",
    ".org",
    ".pptx",
    ".rst",
    ".rtf",
    ".tex",
    ".xlsx",
}
TEXT_EXTENSIONS = {
    ".log",
    ".markdown",
    ".md",
    ".mdx",
    ".sql",
    ".text",
    ".toml",
    ".txt",
    ".yaml",
    ".yml",
}


class ConversionError(RuntimeError):
    pass


@dataclass
class ConversionResult:
    markdown: str
    converter: str
    source_type: str
    warnings: list[str] = field(default_factory=list)


def conversion_capabilities() -> dict[str, bool]:
    return {
        "pandoc": shutil.which("pandoc") is not None,
        "pdftotext": shutil.which("pdftotext") is not None,
        "pdf_python_fallback": importlib.util.find_spec("pypdf") is not None,
        "docx_python": importlib.util.find_spec("docx") is not None,
        "pptx_python": importlib.util.find_spec("pptx") is not None,
        "xlsx_python": importlib.util.find_spec("openpyxl") is not None,
        "pdf_ocr": shutil.which("pdftoppm") is not None and shutil.which("tesseract") is not None,
        "legacy_office": shutil.which("libreoffice") is not None or shutil.which("textutil") is not None,
    }


def convert_document(path: Path) -> ConversionResult:
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTENSIONS:
        return ConversionResult(_read_text(path), "text", suffix.lstrip(".") or "text")
    if suffix == ".pdf":
        return _convert_pdf(path)
    if suffix == ".csv" or suffix == ".tsv":
        return _convert_delimited(path, "\t" if suffix == ".tsv" else ",")
    if suffix == ".json":
        return _convert_json(path)
    if suffix == ".doc":
        return _convert_legacy_doc(path)
    if suffix == ".docx" and importlib.util.find_spec("docx") is not None:
        return _convert_docx(path)
    if suffix == ".pptx" and importlib.util.find_spec("pptx") is not None:
        return _convert_pptx(path)
    if suffix == ".xlsx" and importlib.util.find_spec("openpyxl") is not None:
        return _convert_xlsx(path)
    if suffix == ".ipynb":
        return _convert_notebook(path)
    if suffix in PANDOC_EXTENSIONS:
        return _convert_with_pandoc(path)
    raise ConversionError(f"unsupported document type: {suffix or '<none>'}")


def _read_text(path: Path) -> str:
    try:
        return normalize_markdown(path.read_text(encoding="utf-8"))
    except UnicodeDecodeError:
        return normalize_markdown(path.read_text(encoding="utf-8", errors="replace"))


def _run(command: list[str], *, timeout: int = 180) -> subprocess.CompletedProcess[str]:
    try:
        return subprocess.run(command, capture_output=True, text=True, check=False, timeout=timeout)
    except (OSError, subprocess.TimeoutExpired) as error:
        raise ConversionError(f"converter failed: {command[0]}: {error}") from error


def _convert_with_pandoc(path: Path) -> ConversionResult:
    if shutil.which("pandoc") is None:
        if path.suffix.lower() in {".html", ".htm"}:
            return ConversionResult(_html_fallback(_read_text(path)), "html.parser", "html")
        raise ConversionError(f"pandoc is required to convert {path.suffix.lower()} files")
    result = _run(["pandoc", str(path), "--to=gfm", "--wrap=none"])
    if result.returncode != 0 or not result.stdout.strip():
        detail = result.stderr.strip() or "empty converter output"
        raise ConversionError(f"pandoc could not convert {path.name}: {detail}")
    return ConversionResult(normalize_markdown(result.stdout), "pandoc", path.suffix.lower().lstrip("."))


def _convert_legacy_doc(path: Path) -> ConversionResult:
    with tempfile.TemporaryDirectory() as tempdir:
        temporary = Path(tempdir)
        if shutil.which("libreoffice"):
            result = _run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "docx",
                    "--outdir",
                    str(temporary),
                    str(path),
                ]
            )
            converted = temporary / f"{path.stem}.docx"
            if result.returncode == 0 and converted.exists():
                output = _convert_with_pandoc(converted)
                output.converter = "libreoffice+pandoc"
                return output
        if shutil.which("textutil") and shutil.which("pandoc"):
            converted = temporary / f"{path.stem}.html"
            result = _run(["textutil", "-convert", "html", str(path), "-output", str(converted)])
            if result.returncode == 0 and converted.exists():
                output = _convert_with_pandoc(converted)
                output.converter = "textutil+pandoc"
                return output
    raise ConversionError("legacy .doc conversion requires LibreOffice or textutil plus pandoc")


def _convert_pdf(path: Path) -> ConversionResult:
    pages: list[str] = []
    converter = ""
    warnings: list[str] = []
    if shutil.which("pdftotext"):
        result = _run(["pdftotext", "-layout", "-enc", "UTF-8", str(path), "-"])
        if result.returncode == 0:
            pages = result.stdout.split("\f")
            converter = "pdftotext-layout"
    if not _has_useful_pdf_text(pages):
        fallback = _extract_pdf_with_pypdf(path)
        if _has_useful_pdf_text(fallback):
            pages = fallback
            converter = "pypdf-layout"
    if not _has_useful_pdf_text(pages):
        pages = _ocr_pdf(path)
        if _has_useful_pdf_text(pages):
            converter = "tesseract-ocr"
            warnings.append("PDF had no useful text layer; local OCR fallback was used")
    if not _has_useful_pdf_text(pages):
        raise ConversionError("PDF contains no extractable text and local OCR is unavailable or failed")
    cleaned = _clean_pdf_pages(pages)
    rendered = []
    for page_number, page in enumerate(cleaned, start=1):
        if len(cleaned) > 1:
            rendered.extend([f"## Page {page_number}", "", f"<!-- memento:pdf-page={page_number} -->", ""])
        rendered.extend([page.strip(), ""])
    return ConversionResult(normalize_markdown("\n".join(rendered)), converter, "pdf", warnings)


def _extract_pdf_with_pypdf(path: Path) -> list[str]:
    if importlib.util.find_spec("pypdf") is None:
        return []
    try:
        from pypdf import PdfReader

        reader = PdfReader(path)
        return [
            page.extract_text(extraction_mode="layout", layout_mode_space_vertically=False) or ""
            for page in reader.pages
        ]
    except Exception:
        return []


def _ocr_pdf(path: Path) -> list[str]:
    if shutil.which("pdftoppm") is None or shutil.which("tesseract") is None:
        return []
    with tempfile.TemporaryDirectory() as tempdir:
        prefix = Path(tempdir) / "page"
        rendered = _run(["pdftoppm", "-png", "-r", "200", str(path), str(prefix)], timeout=600)
        if rendered.returncode != 0:
            return []
        languages = _tesseract_languages()
        preferred = [value for value in ("por", "eng") if value in languages]
        language = "+".join(preferred) if preferred else None
        pages = []
        for image in sorted(Path(tempdir).glob("page-*.png")):
            command = ["tesseract", str(image), "stdout", "--psm", "1"]
            if language:
                command.extend(["-l", language])
            result = _run(command, timeout=300)
            pages.append(result.stdout if result.returncode == 0 else "")
        return pages


def _tesseract_languages() -> set[str]:
    result = _run(["tesseract", "--list-langs"])
    if result.returncode != 0:
        return set()
    return {line.strip() for line in result.stdout.splitlines()[1:] if line.strip()}


def _has_useful_pdf_text(pages: list[str]) -> bool:
    alphanumeric = sum(char.isalnum() for page in pages for char in page)
    return alphanumeric >= max(24, len(pages) * 8)


def _clean_pdf_pages(pages: list[str]) -> list[str]:
    normalized = [_dehyphenate_page(page) for page in pages if page.strip()]
    if len(normalized) < 3:
        return normalized
    edges: list[tuple[str, str]] = []
    for page in normalized:
        lines = [line.strip() for line in page.splitlines() if line.strip()]
        edges.append((lines[0] if lines else "", lines[-1] if lines else ""))
    threshold = max(2, (len(normalized) * 3 + 4) // 5)
    repeated = {
        value
        for value, count in Counter(value for pair in edges for value in pair if len(value) >= 3).items()
        if count >= threshold
    }
    cleaned = []
    for page in normalized:
        lines = page.splitlines()
        while lines and lines[0].strip() in repeated:
            lines.pop(0)
        while lines and lines[-1].strip() in repeated:
            lines.pop()
        cleaned.append("\n".join(lines).strip())
    return cleaned


def _dehyphenate_page(page: str) -> str:
    lines = page.replace("\r", "").splitlines()
    output: list[str] = []
    index = 0
    while index < len(lines):
        current = lines[index].rstrip()
        if index + 1 < len(lines) and re.search(r"[A-Za-zÀ-ÿ]-$", current):
            following = lines[index + 1].lstrip()
            if re.match(r"[a-zà-ÿ]", following):
                output.append(current[:-1] + following)
                index += 2
                continue
        output.append(current)
        index += 1
    return "\n".join(output)


def _convert_delimited(path: Path, delimiter: str) -> ConversionResult:
    content = path.read_text(encoding="utf-8-sig", errors="replace")
    rows = list(csv.reader(io.StringIO(content), delimiter=delimiter))
    if not rows:
        raise ConversionError(f"{path.name} has no rows")
    width = max(len(row) for row in rows)
    header = _pad_row(rows[0], width)
    lines = [_markdown_row(header), _markdown_row(["---"] * width)]
    lines.extend(_markdown_row(_pad_row(row, width)) for row in rows[1:])
    return ConversionResult(normalize_markdown("\n".join(lines)), "python-csv", path.suffix.lstrip("."))


def _convert_docx(path: Path) -> ConversionResult:
    from docx import Document

    document = Document(path)
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = ((paragraph.style.name if paragraph.style else "") or "").lower()
        if style.startswith("heading"):
            match = re.search(r"(\d+)", style)
            level = min(int(match.group(1)), 6) if match else 2
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)
    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        rendered = _render_table(rows)
        if rendered:
            parts.append("\n".join(rendered))
    if not parts:
        raise ConversionError(f"{path.name} contains no extractable text")
    return ConversionResult(normalize_markdown("\n\n".join(parts)), "python-docx", "docx")


def _convert_pptx(path: Path) -> ConversionResult:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts = [f"## Slide {slide_number}"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()
                )
                if text:
                    slide_parts.append(text)
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                rendered = _render_table(rows)
                if rendered:
                    slide_parts.append("\n".join(rendered))
        if len(slide_parts) > 1:
            parts.append("\n\n".join(slide_parts))
    if not parts:
        raise ConversionError(f"{path.name} contains no extractable text")
    return ConversionResult(normalize_markdown("\n\n".join(parts)), "python-pptx", "pptx")


def _convert_xlsx(path: Path) -> ConversionResult:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    warnings: list[str] = []
    max_rows = 10_000
    max_columns = 256
    try:
        for worksheet in workbook.worksheets:
            rows: list[list[str]] = []
            if worksheet.max_column > max_columns:
                warnings.append(f"sheet {worksheet.title}: truncated after {max_columns} columns")
            for row_number, row in enumerate(
                worksheet.iter_rows(max_row=max_rows + 1, max_col=max_columns, values_only=True), start=1
            ):
                if row_number > max_rows:
                    warnings.append(f"sheet {worksheet.title}: truncated after {max_rows} rows")
                    break
                values = ["" if value is None else str(value) for value in row[:max_columns]]
                while values and not values[-1]:
                    values.pop()
                if values:
                    rows.append(values)
            if rows:
                parts.append(f"## Sheet: {worksheet.title}\n\n" + "\n".join(_render_table(rows)))
    finally:
        workbook.close()
    if not parts:
        raise ConversionError(f"{path.name} contains no extractable cells")
    return ConversionResult(normalize_markdown("\n\n".join(parts)), "openpyxl", "xlsx", warnings)


def _convert_notebook(path: Path) -> ConversionResult:
    notebook = json.loads(path.read_text(encoding="utf-8-sig"))
    parts: list[str] = []
    for cell in notebook.get("cells", []):
        source = "".join(cell.get("source", [])).strip()
        if not source:
            continue
        if cell.get("cell_type") == "code":
            parts.append(f"```python\n{source}\n```")
        else:
            parts.append(source)
    if not parts:
        raise ConversionError(f"{path.name} contains no readable cells")
    return ConversionResult(normalize_markdown("\n\n".join(parts)), "python-json", "ipynb")


def _render_table(rows: list[list[object]]) -> list[str]:
    if not rows:
        return []
    string_rows = [["" if value is None else str(value) for value in row] for row in rows]
    width = max(len(row) for row in string_rows)
    header = _pad_row(string_rows[0], width)
    lines = [_markdown_row(header), _markdown_row(["---"] * width)]
    lines.extend(_markdown_row(_pad_row(row, width)) for row in string_rows[1:])
    return lines


def _pad_row(row: list[str], width: int) -> list[str]:
    return row + [""] * (width - len(row))


def _markdown_row(row: list[str]) -> str:
    values = [value.replace("|", "\\|").replace("\n", "<br>").strip() for value in row]
    return f"| {' | '.join(values)} |"


def _convert_json(path: Path) -> ConversionResult:
    value = json.loads(path.read_text(encoding="utf-8-sig"))
    body = f"```json\n{json.dumps(value, indent=2, ensure_ascii=False)}\n```"
    return ConversionResult(normalize_markdown(body), "python-json", "json")


class _TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.href: str | None = None

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"p", "div", "br", "li", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")
        self.href = dict(attrs).get("href") if tag == "a" else self.href

    def handle_endtag(self, tag: str) -> None:
        if tag == "a" and self.href:
            self.parts.append(f" ({self.href})")
            self.href = None
        if tag in {"p", "div", "li", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        self.parts.append(html.unescape(data))


def _html_fallback(content: str) -> str:
    parser = _TextHTMLParser()
    parser.feed(content)
    return normalize_markdown("".join(parser.parts))
